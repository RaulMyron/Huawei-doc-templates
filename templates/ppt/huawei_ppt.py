#!/usr/bin/env python3
"""Reusable Huawei Cloud PPT library.

Provides constants, helpers, and slide builders for creating Huawei-branded
slide decks using python-pptx. Extracted from the hotline training guide
monolith into a standalone, importable module.

Usage:
    from huawei_ppt import new_deck, text_box, add_table, callout, ...

Brand colors are locked (AGENTS.md L9). Callout names are locked (L3):
warning, tip, infobox.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from lxml import etree
import os, subprocess, zipfile

# Default template path (auto-detected relative to this file)
_DEFAULT_TEMPLATE = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    'common-assets', 'huawei-template.pptx'
)

# ── Slide dimensions ────────────────────────────────────────────────
SLIDE_W = 13.3   # inches
SLIDE_H = 7.5    # inches

# ── Layout constants ────────────────────────────────────────────────
LEFT_MARGIN = 0.8
RIGHT_MARGIN = 0.8
CONTENT_WIDTH = SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN  # 11.7
TOP_CONTENT = 1.5   # top position for content below title
CENTER_X = SLIDE_W / 2  # 6.65

# ── Brand colors (locked — AGENTS.md L9) ────────────────────────────
RED = RGBColor(0xC7, 0x00, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE = RGBColor(0xE8, 0xE8, 0xE8)
GRAY_BG = RGBColor(0xF6, 0xF8, 0xFA)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK = RGBColor(0x1F, 0x23, 0x28)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_DARK = RGBColor(0x4A, 0x4A, 0x4A)

# ── Callout colors (corporate-muted — AGENTS.md L3 names) ──────────
# Toned down from Bootstrap alert palette to muted, enterprise tones.
AMBER_BG = RGBColor(0xFD, 0xF8, 0xEE)  # subtle warm white
AMBER_FG = RGBColor(0x8A, 0x6D, 0x00)  # muted amber (icon/label)
AMBER_BD = RGBColor(0xD4, 0xA7, 0x2C)  # muted gold border
GREEN_BG = RGBColor(0xED, 0xF6, 0xED)  # subtle green tint
GREEN_FG = RGBColor(0x2E, 0x6B, 0x2E)  # muted green (icon/label)
GREEN_BD = RGBColor(0x5B, 0xA8, 0x5B)  # muted sage border
BLUE_BG = RGBColor(0xED, 0xF3, 0xF9)  # subtle blue tint
BLUE_FG = RGBColor(0x1B, 0x51, 0x70)  # muted blue (icon/label)
BLUE_BD = RGBColor(0x4A, 0x8B, 0xB5)  # muted slate border


# ── Helpers ─────────────────────────────────────────────────────────

def get_layouts(prs):
    """Return a dict mapping layout name → layout object for the presentation."""
    layouts = {}
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            layouts[layout.name] = layout
    return layouts


def remove_all_slides(prs):
    """Remove all slides from the presentation (start from a clean deck)."""
    sldIdLst = prs._element.sldIdLst
    for child in list(sldIdLst):
        sldIdLst.remove(child)


def clean_zip(filepath):
    """Rewrite a ZIP/PPTX with deterministic ordering (no timestamp noise)."""
    with zipfile.ZipFile(filepath, 'r') as zin:
        entries = {}
        for info in zin.infolist():
            entries[info.filename] = (info, zin.read(info.filename))
    with zipfile.ZipFile(filepath, 'w', zipfile.ZIP_DEFLATED) as zout:
        for fn, (info, data) in entries.items():
            zout.writestr(info, data)


def new_deck(template_path=None):
    """Create a new presentation from the Huawei template, with all slides removed.

    Args:
        template_path: Path to the .pptx template. Defaults to the bundled
            ``common-assets/huawei-template.pptx``.

    Returns:
        Tuple of (presentation, layouts_dict).
    """
    if template_path is None:
        template_path = _DEFAULT_TEMPLATE
    prs = Presentation(template_path)
    remove_all_slides(prs)
    return prs, get_layouts(prs)


def add_slide(prs, layouts, layout_name):
    """Add a slide using the named layout.

    Args:
        prs: Presentation object.
        layouts: Layouts dict from :func:`get_layouts`.
        layout_name: Name of the slide layout (e.g. ``"Title Only"``).

    Returns:
        The new slide.
    """
    return prs.slides.add_slide(layouts[layout_name])


def _emu(v):
    """Convert to EMU. Accepts inches (float) or already-EMU (int/Emu) values."""
    if isinstance(v, int):
        return v  # Already EMU (Inches/Emu are int subclasses)
    return Inches(v)


def text_box(slide, text, left, top, width, height,
             size=14, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box with word wrap. Supports multi-line text via ``\\n``.

    Args:
        slide: Slide to add the text box to.
        text: Text content (newlines create separate paragraphs).
        left, top, width, height: Position and size (inches or EMU).
        size: Font size in points (default 14).
        color: Font color (default DARK).
        bold: Whether text is bold (default False).
        align: Paragraph alignment (default LEFT).

    Returns:
        The text box shape.
    """
    tb = slide.shapes.add_textbox(_emu(left), _emu(top), _emu(width), _emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if i > 0:
            p.space_before = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return tb


def set_title(slide, title, color=RED, size=24):
    """Set the slide title placeholder text and style.

    Args:
        slide: Slide whose title placeholder to set.
        title: Title text.
        color: Font color (default RED).
        size: Font size in points (default 24).
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size)
                    r.font.color.rgb = color
                    r.font.bold = True
            return


# ── Table ───────────────────────────────────────────────────────────

def add_table(slide, headers, rows,
              left=None, top=TOP_CONTENT + 0.3,
              col_widths=None):
    """Add a styled table with Huawei red header, alternating rows, cell padding.

    The table is horizontally centered on the slide by default.

    Args:
        slide: Slide to add the table to.
        headers: List of header cell strings.
        rows: List of row lists (each row is a list of cell strings).
        left: Left position in inches (default: centered on slide).
        top: Top position in inches (default TOP_CONTENT + 0.3).
        col_widths: Optional list of column widths (EMU or inches).

    Returns:
        The table shape.
    """
    nr, nc = len(rows) + 1, len(headers)
    if col_widths:
        width = sum(col_widths)
    else:
        width = Inches(CONTENT_WIDTH)
    # Center the table on the slide if left is not specified
    if left is None:
        width_in = width / 914400  # EMU to inches
        left = (SLIDE_W - width_in) / 2
    height = Inches(0.4 * len(rows) + 0.6)
    ts = slide.shapes.add_table(nr, nc, _emu(left), _emu(top), width, height)
    tbl = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w

    # Style header row
    for ci, h in enumerate(headers):
        c = tbl.cell(0, ci)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = RED
        c.margin_left = Inches(0.12)
        c.margin_right = Inches(0.12)
        c.margin_top = Inches(0.08)
        c.margin_bottom = Inches(0.08)
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(13)
                r.font.bold = True
                r.font.color.rgb = WHITE
        c.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Style data rows
    for ri, row in enumerate(rows):
        bg = GRAY_BG if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            c.text = val
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            c.margin_left = Inches(0.12)
            c.margin_right = Inches(0.12)
            c.margin_top = Inches(0.06)
            c.margin_bottom = Inches(0.06)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.color.rgb = DARK
                    if ci == 0:
                        r.font.bold = True
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
    return ts


def table_bottom(ts):
    """Return the bottom Y position of a table shape in inches.

    Useful for positioning callouts below a table.

    Args:
        ts: Table shape returned by :func:`add_table`.

    Returns:
        Bottom Y position in inches.
    """
    return (ts.top + ts.height) / 914400


# ── Callout boxes ───────────────────────────────────────────────────

def callout(slide, kind, text, left=LEFT_MARGIN, top=None, width=CONTENT_WIDTH):
    """Add a callout box with colored background, border, and icon prefix.

    Callout names are locked (AGENTS.md L3): ``warning``, ``tip``,
    ``infobox``.

    Args:
        slide: Slide to add the callout to.
        kind: One of ``'warning'``, ``'tip'``, ``'infobox'``.
        text: Callout text content.
        left: Left position in inches (default LEFT_MARGIN).
        top: Top position in inches (default 4.5 if None).
        width: Width in inches (default CONTENT_WIDTH).

    Returns:
        The text box shape.
    """
    if top is None:
        top = 4.5
    colors = {
        'warning': (AMBER_BG, AMBER_FG, AMBER_BD, '\u26a0'),
        'tip':     (GREEN_BG, GREEN_FG, GREEN_BD, '\u2713'),
        'infobox': (BLUE_BG,  BLUE_FG,  BLUE_BD,  '\u2139'),
    }
    bg, fg, bd, icon = colors.get(kind, colors['infobox'])
    tb = slide.shapes.add_textbox(_emu(left), _emu(top), _emu(width), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    # Set background fill via XML (no border — just the fill)
    spPr = tb._element.spPr
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    clr = etree.SubElement(sf, qn('a:srgbClr'))
    clr.set('val', f'{bg[0]:02X}{bg[1]:02X}{bg[2]:02X}')
    # Add icon prefix + text
    p = tf.paragraphs[0]
    # Icon run (black, bold)
    r_icon = p.add_run()
    r_icon.text = f"{icon}  "
    r_icon.font.size = Pt(14)
    r_icon.font.color.rgb = DARK
    r_icon.font.bold = True
    # Text run (dark, regular — corporate, readable)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(12)
    r.font.color.rgb = DARK
    r.font.bold = False
    return tb


# ── Slide layout builders ───────────────────────────────────────────

def title_slide(prs, layouts, module, subtitle, tag):
    """Title slide — NO "HUAWEI CLOUD" text (template has logo), all text DARK.

    Args:
        prs: Presentation object.
        layouts: Layout dict.
        module: Main title text (large, black, centered).
        subtitle: Subtitle text (black, centered).
        tag: Tag line text (medium gray, centered).
    """
    s = add_slide(prs, layouts, "Title Slide 1")
    # Module name - centered, large, black (on light area of gradient)
    text_box(s, module, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.5),
             36, DARK, True, PP_ALIGN.CENTER)
    # Subtitle - centered, white (on darker area of gradient)
    text_box(s, subtitle, Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.8),
             18, WHITE, False, PP_ALIGN.CENTER)
    # Tag - centered, near-white (on darkest area of gradient)
    text_box(s, tag, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.5),
             14, NEAR_WHITE, False, PP_ALIGN.CENTER)


def chapter_slide(prs, layouts, title, subtitle=""):
    """Chapter divider — large red title centered, subtitle below in lighter color.

    Args:
        prs: Presentation object.
        layouts: Layouts dict.
        title: Chapter title (40pt, bold, red, centered).
        subtitle: Optional subtitle (18pt, medium gray, centered).
    """
    s = add_slide(prs, layouts, "Title Only")
    set_title(s, "")
    # Title: 40pt, bold, red, centered both H and V
    text_box(s, title, Inches(LEFT_MARGIN), Inches(2.2), Inches(CONTENT_WIDTH), Inches(1.5),
             40, RED, True, PP_ALIGN.CENTER)
    # Subtitle: 18pt, centered, medium gray
    if subtitle:
        text_box(s, subtitle, Inches(LEFT_MARGIN), Inches(3.8), Inches(CONTENT_WIDTH), Inches(1.0),
                 18, MED_GRAY, False, PP_ALIGN.CENTER)


def content_slide(prs, layouts, title):
    """Content slide with title set.

    Args:
        prs: Presentation object.
        layouts: Layouts dict.
        title: Slide title.

    Returns:
        The new slide (add content to it).
    """
    s = add_slide(prs, layouts, "Title Only")
    set_title(s, title)
    return s


def last_slide(prs, layouts):
    """Add the closing/last slide.

    Args:
        prs: Presentation object.
        layouts: Layouts dict.
    """
    add_slide(prs, layouts, "Last Page")


# ── Save / export ───────────────────────────────────────────────────

def save_deck(prs, filepath):
    """Save the presentation and clean the ZIP for deterministic output.

    Args:
        prs: Presentation object.
        filepath: Output path for the .pptx file.

    Returns:
        The absolute path to the saved file.
    """
    filepath = os.path.abspath(filepath)
    prs.save(filepath)
    clean_zip(filepath)
    return filepath


def to_pdf(pptx_path):
    """Convert a .pptx file to PDF using LibreOffice.

    Args:
        pptx_path: Path to the .pptx file.

    Returns:
        The path to the generated .pdf file (same directory, same stem).

    Raises:
        FileNotFoundError: If LibreOffice is not installed.
        subprocess.CalledProcessError: If conversion fails.
    """
    pptx_path = os.path.abspath(pptx_path)
    outdir = os.path.dirname(pptx_path)
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', outdir, pptx_path],
        capture_output=True, text=True, check=True
    )
    pdf_path = os.path.splitext(pptx_path)[0] + '.pdf'
    return pdf_path
